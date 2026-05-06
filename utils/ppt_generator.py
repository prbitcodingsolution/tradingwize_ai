"""
Automated PPT Generator for Stock Analysis
Fetches data from database and generates professional PowerPoint presentations
NOW WITH BILINGUAL SUPPORT (English + Hindi)
GENERATES SEPARATE PPT FILES FOR EACH LANGUAGE
"""

import os
import json
from typing import Dict, List, Any, Optional, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import sys
import platform
import subprocess

# Add parent directory to path for imports
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from database_utility.database import StockDatabase
from dotenv import load_dotenv
import openai

load_dotenv()


class  StockPPTGenerator:
    """Generate professional PowerPoint presentations from stock analysis data with bilingual support"""
    
    def __init__(self):
        """Initialize PPT generator with OpenRouter LLM"""
        self.openrouter_api_key = os.getenv('OPENROUTER_API_KEY')
        self.openrouter_base_url = os.getenv('OPENROUTER_BASE_URL', 'https://openrouter.ai/api/v1')
        
        # Configure OpenAI client to use OpenRouter
        self.client = openai.OpenAI(
            api_key=self.openrouter_api_key,
            base_url=self.openrouter_base_url
        )
        
        self.db = StockDatabase()
    
    def translate_to_hindi(self, text_or_list: Union[str, List[str]]) -> Union[str, List[str]]:
        """
        Translate English text to Hindi using LLM.
        For lists: translates each item individually to avoid delimiter/mismatch issues.
        For strings: translates the full text as one block.
        """
        _sys_msg = (
            "You are a professional Hindi translator specializing in financial content. "
            "Translate English to PURE MODERN HINDI using Devanagari script. "
            "DO NOT create Hinglish (mixing English words in Hindi sentences). "
            "Translate ALL common words to Hindi. "
            "Only keep stock tickers (TCS.NS), financial abbreviations (EBITDA, CAGR, EPS, PE, ROE, ROCE), "
            "currency symbols (₹, $, Cr), numbers, and proper names in English. "
            "Return ONLY the Hindi translation, nothing else."
        )

        _translation_rules = """Translate to PURE MODERN HINDI (Devanagari script) for an investor presentation.

TRANSLATE these words to Hindi:
company→कंपनी, market→बाजार, business→व्यवसाय, growth→वृद्धि, revenue→राजस्व, profit→लाभ,
services→सेवाएं, customers→ग्राहक, industry→उद्योग, technology→प्रौद्योगिकी, global→वैश्विक,
leading→अग्रणी, provider→प्रदाता, consulting→परामर्श, strong/robust→मजबूत, financial→वित्तीय,
performance→प्रदर्शन, positive→सकारात्मक, outlook→दृष्टिकोण, investment→निवेश, investor→निवेशक,
opportunity→अवसर, dividend→लाभांश, stability→स्थिरता, potential→क्षमता, future→भविष्य

KEEP in English: stock tickers, EBITDA/CAGR/EPS/PE/ROE/ROCE, ₹/$/Cr, numbers, proper names, IT/AI/CEO

"""

        try:
            # Handle list input — translate each item individually
            if isinstance(text_or_list, list):
                if not text_or_list:
                    return []

                # Filter out None/empty items, track indices
                hindi_results = list(text_or_list)  # start with English copy
                for i, item in enumerate(text_or_list):
                    if not item or not str(item).strip():
                        continue
                    try:
                        response = self.client.chat.completions.create(
                            model="openai/gpt-oss-120b",
                            messages=[
                                {"role": "system", "content": _sys_msg},
                                {"role": "user", "content": _translation_rules + f"Translate this:\n{item}"}
                            ],
                            temperature=0.3,
                            max_tokens=500
                        )
                        result = response.choices[0].message.content
                        if result and result.strip():
                            hindi_results[i] = result.strip()
                    except Exception as item_err:
                        print(f"⚠️ Translation failed for item {i}: {item_err}")
                        # Keep English for this item

                return hindi_results

            # Handle string input
            else:
                if not text_or_list or not str(text_or_list).strip():
                    return ""

                response = self.client.chat.completions.create(
                    model="openai/gpt-oss-120b",
                    messages=[
                        {"role": "system", "content": _sys_msg},
                        {"role": "user", "content": _translation_rules + f"Translate this:\n{text_or_list}"}
                    ],
                    temperature=0.3,
                    max_tokens=1500
                )
                result = response.choices[0].message.content
                return result.strip() if result else text_or_list

        except Exception as e:
            print(f"⚠️ Translation error: {e}. Using English as fallback.")
            return text_or_list  # Fallback to English
    
    def fetch_stock_data(self, stock_symbol: str) -> Optional[Dict[str, Any]]:
        """
        Fetch latest stock data from database
        
        Args:
            stock_symbol: Stock ticker symbol (e.g., 'RELIANCE.NS')
        
        Returns:
            Dict with stock data or None if not found
        """
        try:
            if not self.db.connect():
                print("❌ Failed to connect to database")
                return None
            
            # Get latest analysis
            data = self.db.get_latest_analysis(stock_symbol)
            
            if not data:
                print(f"❌ No data found for {stock_symbol}")
                return None
            
            # Structure the data for LLM.
            # The client disabled the market-sentiment pipeline, so we
            # intentionally do NOT pull `market_senti` — the PPT now builds
            # from three DB columns: analyzed_response (fundamentals),
            # future_senti (future outlook block), and finrobot_response
            # (FinRobot deep-analysis memo). market_senti is still exposed
            # as "Not used — pipeline disabled" so any legacy prompt that
            # referenced the key fails loudly instead of silently reading
            # empty values.
            structured_data = {
                "stock_name": data['stock_name'],
                "stock_symbol": data['stock_symbol'],
                "fundamentals": data['analyzed_response'],
                "sentiment": {
                    "current_market_sentiment": "Not used — pipeline disabled",
                    "current_sentiment_status": "disabled",
                    "future_outlook": data.get('future_senti', 'Not available'),
                    "future_sentiment_status": data.get('future_senti_status', 'neutral'),
                },
                "finrobot": {
                    "response": data.get('finrobot_response', '') or '',
                    "recommendation": data.get('finrobot_recommendation', '') or '',
                    "score": data.get('finrobot_score'),
                },
                "tech_analysis": data.get('tech_analysis', {}),
                "selection": data.get('selection', False),
                "analyzed_at": str(data.get('analyzed_at', '')),
            }

            self.db.disconnect()
            return structured_data
            
        except Exception as e:
            print(f"❌ Error fetching stock data: {e}")
            if self.db.conn:
                self.db.disconnect()
            return None
    
    def generate_slide_structure(self, stock_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Use LLM to generate PPT slide structure with comprehensive ENGLISH content
        Then translate to Hindi for bilingual support
        
        Args:
            stock_data: Structured stock data from database
        
        Returns:
            Dict with slide structure including both English and Hindi content
        """
        try:
            print("🤖 Generating comprehensive slide structure using LLM (English)...")
            
            # Safely extract the three content sources the PPT is built
            # from now that market-sentiment is disabled:
            #   1. fundamentals (analyzed_response)
            #   2. future_senti (future outlook block)
            #   3. finrobot_response (FinRobot deep-analysis memo)
            future_outlook = stock_data['sentiment'].get('future_outlook')
            future_sentiment_status = stock_data['sentiment'].get('future_sentiment_status', 'neutral')

            if future_outlook and future_outlook not in ("Not available", ""):
                future_outlook_text = future_outlook[:3000]
            else:
                future_outlook_text = (
                    'Future outlook not available — re-run the main analysis on '
                    'the Data Dashboard so the future_senti column is populated.'
                )

            finrobot_block = stock_data.get('finrobot', {}) or {}
            finrobot_text_raw = (finrobot_block.get('response') or '').strip()
            finrobot_rec = finrobot_block.get('recommendation') or 'Not yet available'
            finrobot_score = finrobot_block.get('score')
            finrobot_score_str = (
                f"{float(finrobot_score):.1f}/100"
                if finrobot_score is not None else "N/A"
            )
            if finrobot_text_raw:
                # Keep plenty of room — this is the single richest narrative
                # source and carries the reasoning chain + bull/bear case +
                # catalysts the recommendation slide should anchor on.
                finrobot_text = finrobot_text_raw[:5000]
            else:
                finrobot_text = (
                    'FinRobot deep analysis not saved yet — click "🚀 Run Deep '
                    'Analysis" on the Deep Analysis → FinRobot Agent tab.'
                )

            # Create enhanced prompt for LLM (ENGLISH ONLY - we'll translate after)
            prompt = f"""You are a senior financial analyst creating a comprehensive, professional stock analysis presentation for investors.

Based on the following detailed stock data, generate a rich PPT structure with 10-12 slides that includes BOTH bullet points AND paragraph-style content.

Stock Data:
Stock Name: {stock_data['stock_name']}
Stock Symbol: {stock_data['stock_symbol']}
Selection Status: {'SELECTED' if stock_data['selection'] else 'NOT SELECTED'}
Future Outlook: {future_sentiment_status}
FinRobot Recommendation: {finrobot_rec} (Score: {finrobot_score_str})

Fundamental Analysis (analyzed_response):
{stock_data['fundamentals'][:3000]}

Future Outlook (future_senti):
{future_outlook_text}

FinRobot Deep-Analysis Memo (finrobot_response):
{finrobot_text}

Technical Metrics:
{json.dumps(stock_data['tech_analysis'], indent=2)}

Generate a JSON structure with the following format:
{{
  "slides": [
    {{
      "title": "Slide Title",
      "type": "bullets" or "paragraph" or "mixed",
      "content": ["Point 1", "Point 2", ...] for bullets,
      "paragraph": "Full paragraph text..." for paragraph slides,
      "summary": "Brief summary of the slide" (optional)
    }}
  ]
}}

CRITICAL REQUIREMENTS:

1. Create 10-12 comprehensive slides:
   - Slide 1: Title slide (bullets) - Company name, symbol, date
   - Slide 2: Executive Summary (paragraph) - 150-200 word overview of the company and investment thesis, leaning on the FinRobot memo's executive-summary + investment-thesis sections
   - Slide 3: Company Overview (mixed) - 2-3 bullets + 100-word paragraph about business model (from analyzed_response)
   - Slide 4: Financial Performance (bullets) - 6-8 detailed financial metrics with actual numbers (from analyzed_response + Technical Metrics JSON)
   - Slide 5: Financial Analysis (paragraph) - 150-200 word analysis of financial health and trends, incorporating the FinRobot valuation/financial-health/growth commentary
   - Slide 6: Technical Analysis (bullets) - 5-7 key technical indicators with values
   - Slide 7: Market Position & Competitive Landscape (mixed) - 3-4 bullets + 100-word paragraph; pull peer-comparison + moat from the FinRobot memo when present
   - Slide 8: Future Outlook & Growth Drivers (paragraph) - 150-200 word analysis of future prospects (from future_senti block)
   - Slide 9: FinRobot Deep Analysis (mixed) - show recommendation + score, a 120-word extract of the chain-of-thought / investment thesis, and 3-4 bullet catalysts from the FinRobot memo
   - Slide 10: Bull Case vs Bear Case (mixed) - 3-4 bullets for each side; if the FinRobot memo has bull_case/bear_case lists, USE them; otherwise synthesise from fundamentals + future_senti
   - Slide 11: Risk Factors (bullets) - 6-8 specific risks with details (merge fundamentals key_risks + future_senti risk_factors + FinRobot bear_case)
   - Slide 12: Investment Recommendation & Next Steps (paragraph) - 120-180 words concluding with the FinRobot recommendation, confidence, time horizon, and price-level guidance when available

2. Content Guidelines:
   - For BULLET slides: 5-8 detailed points (not just 3-5)
   - For PARAGRAPH slides: 100-200 words of flowing narrative
   - For MIXED slides: 3-4 bullets + 80-120 word paragraph
   - Include ALL specific numbers, percentages, and metrics from the data
   - Use professional, investor-friendly language
   - Make content actionable and insightful

3. Paragraph Content Should:
   - Tell a story about the company/analysis
   - Connect different data points
   - Provide context and interpretation
   - Be engaging and easy to read
   - Include specific examples and numbers

4. Bullet Points Should:
   - Be specific and data-driven
   - Include actual numbers and metrics
   - Be concise but informative (10-20 words each)
   - Provide actionable insights

5. Extract and use:
   - All financial metrics (revenue, profit, margins, PE ratio, etc.)
   - All technical indicators (price, market cap, 52-week high/low, etc.)
   - Sentiment scores and analysis
   - Growth trends and forecasts
   - Competitive advantages
   - Risk factors

Return ONLY the JSON structure, no additional text or markdown."""

            # Call LLM with a big-enough token budget for 12 detailed
            # slides. The old 4000-token cap was truncating the JSON
            # response mid-way, which triggered JSONDecodeError and
            # dropped callers into the 4-slide fallback — exactly the
            # bug the client reported.
            response = self.client.chat.completions.create(
                model="openai/gpt-oss-120b",
                messages=[
                    {"role": "system", "content": "You are a senior financial analyst expert at creating comprehensive, professional stock analysis presentations with rich content. Generate detailed slides with both bullet points and paragraph-style content. Return only valid JSON — never truncate mid-object."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.6,
                max_tokens=12000,
                timeout=90.0,
            )

            # Extract and parse response
            llm_output = response.choices[0].message.content.strip()

            # Remove markdown code blocks if present
            if llm_output.startswith("```json"):
                llm_output = llm_output[7:]
            if llm_output.startswith("```"):
                llm_output = llm_output[3:]
            if llm_output.endswith("```"):
                llm_output = llm_output[:-3]

            llm_output = llm_output.strip()

            # Parse JSON with a recovery pass for truncated responses.
            slide_structure = self._parse_slide_json(llm_output)

            _slide_count = len(slide_structure.get('slides', []))
            print(f"✅ Generated {_slide_count} comprehensive slides (English)")
            if _slide_count < 10:
                # Enough slides were lost to the token cap or JSON recovery
                # pass to be worth reporting — the final PPT will still
                # render but may come in short of the 12-slide target.
                print(
                    f"⚠️ Only {_slide_count} slides generated; expected 10-12. "
                    "Consider increasing max_tokens or trimming the prompt context."
                )
            
            # NOW TRANSLATE TO HINDI
            print("🌐 Translating content to Hindi...")
            
            for i, slide in enumerate(slide_structure.get('slides', []), 1):
                slide_type = slide.get('type', 'bullets')
                
                # Handle None content by converting to empty list
                content = slide.get('content', [])
                if content is None:
                    content = []
                    slide['content'] = []
                
                # Translate based on slide type
                if slide_type == 'bullets':
                    # Translate bullet points
                    if content:
                        print(f"   Translating slide {i}: {slide.get('title', 'Untitled')} ({len(content)} bullets)")
                        hindi_content = self.translate_to_hindi(content)
                        slide['content_eng'] = content
                        slide['content_hindi'] = hindi_content
                        # Keep 'content' for backward compatibility
                        slide['content'] = content
                    else:
                        slide['content_eng'] = []
                        slide['content_hindi'] = []
                
                elif slide_type == 'paragraph':
                    # Translate paragraph
                    paragraph = slide.get('paragraph', '')
                    if paragraph:
                        print(f"   Translating slide {i}: {slide.get('title', 'Untitled')} (paragraph)")
                        hindi_paragraph = self.translate_to_hindi(paragraph)
                        slide['paragraph_eng'] = paragraph
                        slide['paragraph_hindi'] = hindi_paragraph
                        # Keep 'paragraph' for backward compatibility
                        slide['paragraph'] = paragraph
                    else:
                        slide['paragraph_eng'] = ''
                        slide['paragraph_hindi'] = ''
                
                elif slide_type == 'mixed':
                    # Translate both bullets and paragraph
                    if content:
                        print(f"   Translating slide {i}: {slide.get('title', 'Untitled')} (mixed: {len(content)} bullets + paragraph)")
                        hindi_content = self.translate_to_hindi(content)
                        slide['content_eng'] = content
                        slide['content_hindi'] = hindi_content
                        slide['content'] = content
                    else:
                        slide['content_eng'] = []
                        slide['content_hindi'] = []
                    
                    paragraph = slide.get('paragraph', '')
                    if paragraph:
                        hindi_paragraph = self.translate_to_hindi(paragraph)
                        slide['paragraph_eng'] = paragraph
                        slide['paragraph_hindi'] = hindi_paragraph
                        slide['paragraph'] = paragraph
                    else:
                        slide['paragraph_eng'] = ''
                        slide['paragraph_hindi'] = ''
            
            print(f"✅ Translation complete! All slides now have English + Hindi content")
            
            # Log slide types
            for i, slide in enumerate(slide_structure.get('slides', []), 1):
                slide_type = slide.get('type', 'bullets')
                content_count = len(slide.get('content_eng', []))
                has_paragraph = 'paragraph_eng' in slide and slide.get('paragraph_eng')
                print(f"   Slide {i}: {slide.get('title', 'Untitled')} ({slide_type}) - {content_count} points" + 
                      (f" + paragraph ({len(slide.get('paragraph_eng', ''))} chars)" if has_paragraph else "") +
                      " [EN + HI]")
            
            # Save JSON to PPT_json folder
            try:
                import os
                from datetime import datetime
                
                # Create PPT_json folder if it doesn't exist
                json_folder = "PPT_json"
                if not os.path.exists(json_folder):
                    os.makedirs(json_folder)
                    print(f"📁 Created folder: {json_folder}")
                
                # Generate filename with stock symbol and timestamp
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                stock_symbol_clean = stock_data['stock_symbol'].replace('.', '_')
                json_filename = f"{stock_symbol_clean}_{timestamp}.json"
                json_filepath = os.path.join(json_folder, json_filename)
                
                # Save the complete structure with metadata
                json_data = {
                    "metadata": {
                        "stock_name": stock_data['stock_name'],
                        "stock_symbol": stock_data['stock_symbol'],
                        "generated_at": datetime.now().isoformat(),
                        "selection_status": stock_data['selection'],
                        "current_sentiment": current_sentiment_status,
                        "future_sentiment": future_sentiment_status,
                        "language_support": "bilingual",  # NEW: Indicate bilingual support
                        "languages": ["english", "hindi"]  # NEW: List supported languages
                    },
                    "slide_structure": slide_structure
                }
                
                with open(json_filepath, 'w', encoding='utf-8') as f:
                    json.dump(json_data, f, indent=2, ensure_ascii=False)
                
                print(f"💾 Saved bilingual PPT JSON to: {json_filepath}")
                
            except Exception as save_error:
                print(f"⚠️ Failed to save JSON file: {save_error}")
                # Continue even if save fails
            
            return slide_structure
            
        except json.JSONDecodeError as e:
            print(f"❌ Failed to parse LLM response as JSON: {e}")
            print(f"LLM Output: {llm_output[:500]}...")
            return self._get_fallback_structure(stock_data)
        except Exception as e:
            print(f"❌ Error generating slide structure: {e}")
            import traceback
            traceback.print_exc()
            return self._get_fallback_structure(stock_data)
    
    def _parse_slide_json(self, raw: str) -> Dict[str, Any]:
        """Parse an LLM JSON response into {slides: [...]}. Tolerates
        common truncation failures from hitting the token cap mid-array
        by recovering up to the last complete slide object.
        """
        # Fast path: fully-formed JSON.
        try:
            data = json.loads(raw)
            if isinstance(data, dict) and "slides" in data:
                return data
            if isinstance(data, list):
                return {"slides": data}
        except json.JSONDecodeError:
            pass

        # Recovery path — find the `"slides": [` array and parse slide
        # objects one at a time, stopping at the first object that fails
        # to parse. This rescues responses that got truncated halfway
        # through a later slide.
        import re as _re

        # Trim anything before the first `[` after "slides"
        _match = _re.search(r'"slides"\s*:\s*\[', raw)
        if not _match:
            raise json.JSONDecodeError("no 'slides' array found in LLM output", raw, 0)

        _cursor = _match.end()
        slides: list[Dict[str, Any]] = []
        while _cursor < len(raw):
            # Skip whitespace and commas
            while _cursor < len(raw) and raw[_cursor] in " \t\r\n,":
                _cursor += 1
            if _cursor >= len(raw) or raw[_cursor] == "]":
                break
            if raw[_cursor] != "{":
                break

            # Scan to the matching closing brace, respecting strings.
            _depth = 0
            _in_string = False
            _escape = False
            _start = _cursor
            while _cursor < len(raw):
                _ch = raw[_cursor]
                if _escape:
                    _escape = False
                elif _ch == "\\":
                    _escape = True
                elif _ch == '"':
                    _in_string = not _in_string
                elif not _in_string:
                    if _ch == "{":
                        _depth += 1
                    elif _ch == "}":
                        _depth -= 1
                        if _depth == 0:
                            _cursor += 1
                            break
                _cursor += 1
            else:
                # Hit EOF before matching brace — truncated slide, drop it.
                break

            _obj_raw = raw[_start:_cursor]
            try:
                slides.append(json.loads(_obj_raw))
            except json.JSONDecodeError:
                # Single malformed slide — stop here, keep what we have.
                break

        if not slides:
            raise json.JSONDecodeError("could not recover any complete slides", raw, 0)

        print(
            f"⚠️ LLM JSON was truncated; recovered {len(slides)} complete slides "
            "via parser fallback."
        )
        return {"slides": slides}

    def _get_fallback_structure(self, stock_data: Dict[str, Any]) -> Dict[str, Any]:
        """Fallback slide structure when the LLM call fails.

        Produces 12 slides with a mix of bullets / paragraph / mixed
        types, built from the actual DB data (analyzed_response,
        future_senti, finrobot_response) so even the fallback path gives
        the user the presentation shape they asked for instead of the
        old 4-slide bullet-only stub.
        """
        name = stock_data.get('stock_name', 'Stock')
        symbol = stock_data.get('stock_symbol', 'N/A')
        selection = bool(stock_data.get('selection', False))
        analyzed_at = str(stock_data.get('analyzed_at', 'N/A'))

        fundamentals = (stock_data.get('fundamentals') or '').strip()
        sentiment_block = stock_data.get('sentiment') or {}
        future_outlook = (sentiment_block.get('future_outlook') or '').strip()
        future_status = sentiment_block.get('future_sentiment_status') or 'neutral'

        finrobot = stock_data.get('finrobot') or {}
        fr_text = (finrobot.get('response') or '').strip()
        fr_rec = finrobot.get('recommendation') or 'Under review'
        fr_score = finrobot.get('score')
        fr_score_str = f"{float(fr_score):.1f}/100" if fr_score is not None else "N/A"

        tech = stock_data.get('tech_analysis') or {}

        def _excerpt(text: str, limit: int = 900) -> str:
            """Condense a potentially long text block into a single
            paragraph-sized extract for a paragraph slide."""
            if not text:
                return "Data not yet available."
            _single = " ".join(text.split())
            return _single[:limit] + ("..." if len(_single) > limit else "")

        def _first_bullets(text: str, n: int = 6, min_len: int = 20) -> list[str]:
            """Pull the first `n` meaningful lines from a long text block
            (e.g. analyzed_response) to seed a bullet slide. Skips
            headers and short/empty lines."""
            if not text:
                return [f"Data not yet available for {name}."]
            out: list[str] = []
            for line in text.splitlines():
                _l = line.strip().lstrip("-•*·#").strip()
                if len(_l) < min_len:
                    continue
                # Skip lines that look like section headers only.
                if _l.endswith(":") and len(_l) < 60:
                    continue
                out.append(_l[:180])
                if len(out) >= n:
                    break
            if not out:
                out = [f"Full analysis available in the {name} report."]
            return out

        tech_bullets = [f"{k}: {v}" for k, v in tech.items() if v not in (None, '', 'N/A')]
        if not tech_bullets:
            tech_bullets = ["Technical metrics not yet populated — re-run main analysis."]

        def _s(type_: str, title: str, *, content: list[str] | None = None,
               paragraph: str | None = None) -> Dict[str, Any]:
            """Assemble a slide dict, mirroring `content`/`paragraph`
            into the `*_eng` / `*_hindi` slots so the PPT renderer and
            language switcher both work. Hindi falls back to English
            text (translation happens only when the LLM path succeeds)."""
            slide: Dict[str, Any] = {"title": title, "type": type_}
            if content is not None:
                slide["content"] = content
                slide["content_eng"] = content
                slide["content_hindi"] = content
            if paragraph is not None:
                slide["paragraph"] = paragraph
                slide["paragraph_eng"] = paragraph
                slide["paragraph_hindi"] = paragraph
            return slide

        return {
            "slides": [
                _s("bullets", f"{name} — Stock Analysis", content=[
                    f"Symbol: {symbol}",
                    f"Analysis Date: {analyzed_at}",
                    f"FinRobot Recommendation: {fr_rec} (Score: {fr_score_str})",
                    f"Selection Status: {'SELECTED' if selection else 'NOT SELECTED'}",
                    "Comprehensive Investment Analysis",
                ]),
                _s("paragraph", "Executive Summary",
                   paragraph=_excerpt(fr_text or fundamentals, 1200)),
                _s("mixed", "Company Overview",
                   content=[
                       f"Company: {name}",
                       f"Ticker: {symbol}",
                       f"Analysed on: {analyzed_at}",
                   ],
                   paragraph=_excerpt(fundamentals, 800)),
                _s("bullets", "Financial Performance",
                   content=_first_bullets(fundamentals, n=8)),
                _s("paragraph", "Financial Analysis",
                   paragraph=_excerpt(fundamentals, 1000)),
                _s("bullets", "Technical Metrics", content=tech_bullets[:8]),
                _s("mixed", "Market Position & Competitive Landscape",
                   content=_first_bullets(fr_text or fundamentals, n=4),
                   paragraph=_excerpt(fr_text or fundamentals, 700)),
                _s("paragraph", "Future Outlook & Growth Drivers",
                   paragraph=_excerpt(future_outlook, 1100)),
                _s("mixed", "FinRobot Deep Analysis",
                   content=[
                       f"Recommendation: {fr_rec}",
                       f"Score: {fr_score_str}",
                       f"Future sentiment status: {future_status}",
                   ],
                   paragraph=_excerpt(fr_text, 1000)),
                _s("bullets", "Bull Case vs Bear Case",
                   content=_first_bullets(fr_text, n=6)),
                _s("bullets", "Risk Factors",
                   content=_first_bullets(fundamentals + "\n" + future_outlook, n=8)),
                _s("paragraph", "Investment Recommendation & Next Steps",
                   paragraph=_excerpt(
                       f"Recommendation: {fr_rec} (Score: {fr_score_str}). "
                       + (fr_text or fundamentals),
                       1000,
                   )),
            ]
        }
    
    def create_presentation(self, slide_structure: Dict[str, Any], stock_data: Dict[str, Any], 
                          output_path: str, language: str = "english") -> bool:
        """
        Create PowerPoint presentation from slide structure in specified language
        
        Args:
            slide_structure: Slide structure from LLM
            stock_data: Original stock data
            output_path: Path to save the PPT file
            language: "english" or "hindi"
        
        Returns:
            bool: True if successful
        """
        try:
            print(f"📊 Creating PowerPoint presentation ({language.upper()})...")
            
            # Create presentation
            prs = Presentation()
            
            # Set slide size to widescreen (16:9)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Process each slide
            for idx, slide_data in enumerate(slide_structure.get('slides', [])):
                slide_type = slide_data.get('type', 'bullets')
                
                if idx == 0:
                    # Title slide
                    self._create_title_slide(prs, slide_data, stock_data, language)
                elif slide_type == 'paragraph':
                    # Paragraph slide
                    self._create_paragraph_slide(prs, slide_data, language)
                elif slide_type == 'mixed':
                    # Mixed slide (bullets + paragraph)
                    self._create_mixed_slide(prs, slide_data, language)
                else:
                    # Standard bullet point slide
                    self._create_content_slide(prs, slide_data, language)
            
            # Save presentation
            prs.save(output_path)
            print(f"✅ {language.upper()} presentation saved: {output_path}")
            return True
            
        except Exception as e:
            print(f"❌ Error creating {language} presentation: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def _create_paragraph_slide(self, prs: Presentation, slide_data: Dict[str, Any], language: str = "english"):
        """Create slide with paragraph content in specified language
        
        Args:
            prs: Presentation object
            slide_data: Slide data with content_eng and content_hindi
            language: "english" or "hindi"
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 51, 102)
        
        # Select content based on language
        if language == "hindi":
            paragraph_text = slide_data.get('paragraph_hindi', slide_data.get('paragraph', ''))
        else:  # english
            paragraph_text = slide_data.get('paragraph_eng', slide_data.get('paragraph', ''))
        
        # Add paragraph content
        if paragraph_text:
            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.text = paragraph_text
            
            # Format paragraph
            for paragraph in content_frame.paragraphs:
                paragraph.font.size = Pt(16)
                paragraph.line_spacing = 1.5
                paragraph.space_after = Pt(12)
                paragraph.alignment = PP_ALIGN.LEFT
    
    def _create_mixed_slide(self, prs: Presentation, slide_data: Dict[str, Any], language: str = "english"):
        """Create slide with both bullets and paragraph in specified language
        
        Args:
            prs: Presentation object
            slide_data: Slide data with content_eng/content_hindi and paragraph_eng/paragraph_hindi
            language: "english" or "hindi"
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 51, 102)
        
        # Select content based on language
        if language == "hindi":
            content_to_show = slide_data.get('content_hindi', slide_data.get('content', []))
            paragraph_text = slide_data.get('paragraph_hindi', slide_data.get('paragraph', ''))
        else:  # english
            content_to_show = slide_data.get('content_eng', slide_data.get('content', []))
            paragraph_text = slide_data.get('paragraph_eng', slide_data.get('paragraph', ''))
        
        # Add bullet points (top half)
        bullet_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.5)
        )
        bullet_frame = bullet_box.text_frame
        bullet_frame.word_wrap = True
        
        for i, point in enumerate(content_to_show):
            if i == 0:
                p = bullet_frame.paragraphs[0]
            else:
                p = bullet_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(16)
            p.space_before = Pt(6)
            p.space_after = Pt(6)
        
        # Add paragraph (bottom half)
        if paragraph_text:
            para_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(4.3), Inches(11.733), Inches(2.7)
            )
            para_frame = para_box.text_frame
            para_frame.word_wrap = True
            para_frame.text = paragraph_text
            
            # Format paragraph
            for paragraph in para_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.line_spacing = 1.4
                paragraph.space_after = Pt(8)
                paragraph.alignment = PP_ALIGN.LEFT
    
    def _create_title_slide(self, prs: Presentation, slide_data: Dict[str, Any], stock_data: Dict[str, Any], language: str = "english"):
        """Create title slide with custom styling in specified language
        
        Args:
            prs: Presentation object
            slide_data: Slide data
            stock_data: Stock information
            language: "english" or "hindi"
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = stock_data.get('stock_name', slide_data['title'])
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Add subtitle
        subtitle_text = f"{stock_data['stock_symbol']} | {datetime.now().strftime('%B %d, %Y')}"
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
    
    def _create_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], language: str = "english"):
        """Create content slide with title and bullet points in specified language
        
        Args:
            prs: Presentation object
            slide_data: Slide data with content_eng and content_hindi
            language: "english" or "hindi"
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout for more control
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 51, 102)
        
        # Select content based on language
        if language == "hindi":
            content_to_show = slide_data.get('content_hindi', slide_data.get('content', []))
        else:  # english
            content_to_show = slide_data.get('content_eng', slide_data.get('content', []))
        
        # Add content with more space for bullets
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Add bullet points
        for i, point in enumerate(content_to_show):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(16)
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.line_spacing = 1.2
    
    def convert_ppt_to_pdf(self, ppt_path: str) -> Optional[str]:
        """
        Convert PowerPoint to PDF
        
        Args:
            ppt_path: Path to the PPT file
        
        Returns:
            str: Path to generated PDF file or None if failed
        """
        try:
            print(f"📄 Converting PPT to PDF...")
            
            # Generate PDF path (same name, different extension)
            pdf_path = ppt_path.replace('.pptx', '.pdf')
            
            # Detect platform and use appropriate conversion method
            system = platform.system()
            
            if system == "Windows":
                # Windows: Use comtypes to convert via PowerPoint COM
                try:
                    import comtypes.client
                    import pythoncom
                    
                    print("🪟 Using Windows PowerPoint COM for conversion...")
                    
                    # Initialize COM for this thread
                    pythoncom.CoInitialize()
                    
                    try:
                        # Initialize PowerPoint
                        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                        powerpoint.Visible = 1
                        
                        # Open presentation
                        abs_ppt_path = os.path.abspath(ppt_path)
                        abs_pdf_path = os.path.abspath(pdf_path)
                        
                        presentation = powerpoint.Presentations.Open(abs_ppt_path, WithWindow=False)
                        
                        # Save as PDF (format type 32 = PDF)
                        presentation.SaveAs(abs_pdf_path, 32)
                        
                        # Close presentation and PowerPoint
                        presentation.Close()
                        powerpoint.Quit()
                        
                        print(f"✅ PDF created: {pdf_path}")
                        return pdf_path
                    
                    finally:
                        # Always uninitialize COM
                        pythoncom.CoUninitialize()
                    
                except ImportError as ie:
                    if 'comtypes' in str(ie):
                        print("⚠️ comtypes not installed. Install with: pip install comtypes")
                    elif 'pythoncom' in str(ie):
                        print("⚠️ pywin32 not installed. Install with: pip install pywin32")
                    else:
                        print(f"⚠️ Import error: {ie}")
                    print("⚠️ Skipping PDF conversion")
                    return None
                except Exception as e:
                    print(f"⚠️ Windows COM conversion failed: {e}")
                    print(f"⚠️ Error type: {type(e).__name__}")
                    return None
            
            elif system == "Linux":
                # Linux: Use LibreOffice
                try:
                    print("🐧 Using LibreOffice for conversion...")
                    
                    # Check if LibreOffice is installed
                    result = subprocess.run(
                        ["libreoffice", "--version"],
                        capture_output=True,
                        text=True,
                        timeout=5
                    )
                    
                    if result.returncode != 0:
                        print("⚠️ LibreOffice not found. Install with: sudo apt-get install libreoffice")
                        return None
                    
                    # Convert using LibreOffice headless mode
                    output_dir = os.path.dirname(os.path.abspath(ppt_path))
                    abs_ppt_path = os.path.abspath(ppt_path)
                    
                    result = subprocess.run(
                        [
                            "libreoffice",
                            "--headless",
                            "--convert-to", "pdf",
                            "--outdir", output_dir,
                            abs_ppt_path
                        ],
                        capture_output=True,
                        text=True,
                        timeout=60
                    )
                    
                    if result.returncode == 0 and os.path.exists(pdf_path):
                        print(f"✅ PDF created: {pdf_path}")
                        return pdf_path
                    else:
                        print(f"⚠️ LibreOffice conversion failed: {result.stderr}")
                        return None
                        
                except FileNotFoundError:
                    print("⚠️ LibreOffice not found. Install with: sudo apt-get install libreoffice")
                    return None
                except subprocess.TimeoutExpired:
                    print("⚠️ LibreOffice conversion timed out")
                    return None
                except Exception as e:
                    print(f"⚠️ Linux conversion failed: {e}")
                    return None
            
            elif system == "Darwin":
                # macOS: Use LibreOffice or Keynote
                print("🍎 macOS detected - PDF conversion not implemented yet")
                print("⚠️ Please install LibreOffice for PDF conversion")
                return None
            
            else:
                print(f"⚠️ Unsupported platform: {system}")
                return None
                
        except Exception as e:
            print(f"❌ Error converting PPT to PDF: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def generate_ppt(self, stock_symbol: str, output_dir: str = "downloads") -> Optional[Dict[str, str]]:
        """
        Main method to generate BILINGUAL PPT for a stock (separate English and Hindi files)
        
        Args:
            stock_symbol: Stock ticker symbol
            output_dir: Directory to save the PPT
        
        Returns:
            dict: {
                'ppt_path_en': str, 
                'ppt_path_hi': str,
                'pdf_path_en': str or None,
                'pdf_path_hi': str or None
            } or None if failed
        """
        try:
            print(f"\n{'='*60}")
            print(f"🚀 Starting BILINGUAL PPT generation for {stock_symbol}")
            print(f"{'='*60}\n")
            
            # Step 1: Fetch data from database
            print("📊 Step 1: Fetching stock data from database...")
            stock_data = self.fetch_stock_data(stock_symbol)
            if not stock_data:
                return None
            print(f"✅ Data fetched for {stock_data['stock_name']}")
            
            # Step 2: Generate slide structure using LLM (with bilingual content)
            print("\n🤖 Step 2: Generating bilingual slide structure...")
            slide_structure = self.generate_slide_structure(stock_data)
            if not slide_structure or 'slides' not in slide_structure:
                print("❌ Failed to generate slide structure")
                return None
            
            # Step 3: Create TWO PowerPoint presentations (English and Hindi)
            print("\n📊 Step 3: Creating PowerPoint presentations (English + Hindi)...")
            
            # Create output directory if it doesn't exist
            os.makedirs(output_dir, exist_ok=True)
            
            # Generate filenames with language suffix
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            base_filename = f"{stock_data['stock_name'].replace(' ', '_')}_{stock_symbol.replace('.', '_')}_{timestamp}"
            
            english_filename = f"{base_filename}_EN.pptx"
            hindi_filename = f"{base_filename}_HI.pptx"
            
            english_path = os.path.join(output_dir, english_filename)
            hindi_path = os.path.join(output_dir, hindi_filename)
            
            # Create English presentation
            print("\n   🇬🇧 Creating English presentation...")
            success_en = self.create_presentation(slide_structure, stock_data, english_path, language="english")
            
            if not success_en:
                print("❌ Failed to create English presentation")
                return None
            
            print(f"   ✅ English PPT created: {english_path}")
            
            # Create Hindi presentation
            print("\n   🇮🇳 Creating Hindi presentation...")
            success_hi = self.create_presentation(slide_structure, stock_data, hindi_path, language="hindi")
            
            if not success_hi:
                print("❌ Failed to create Hindi presentation")
                return None
            
            print(f"   ✅ Hindi PPT created: {hindi_path}")
            
            # Step 4: Convert both to PDF
            print("\n📄 Step 4: Converting to PDF...")
            
            print("   🇬🇧 Converting English PPT to PDF...")
            pdf_path_en = self.convert_ppt_to_pdf(english_path)
            if pdf_path_en:
                print(f"   ✅ English PDF created: {pdf_path_en}")
            else:
                print("   ⚠️ English PDF conversion skipped or failed")
            
            print("   🇮🇳 Converting Hindi PPT to PDF...")
            pdf_path_hi = self.convert_ppt_to_pdf(hindi_path)
            if pdf_path_hi:
                print(f"   ✅ Hindi PDF created: {pdf_path_hi}")
            else:
                print("   ⚠️ Hindi PDF conversion skipped or failed")
            
            print(f"\n{'='*60}")
            print(f"✅ BILINGUAL GENERATION COMPLETE!")
            print(f"📁 English PPT: {english_path}")
            print(f"📁 Hindi PPT: {hindi_path}")
            if pdf_path_en:
                print(f"📁 English PDF: {pdf_path_en}")
            if pdf_path_hi:
                print(f"📁 Hindi PDF: {pdf_path_hi}")
            print(f"{'='*60}\n")
            
            return {
                'ppt_path_en': english_path,
                'ppt_path_hi': hindi_path,
                'pdf_path_en': pdf_path_en,
                'pdf_path_hi': pdf_path_hi
            }
                
        except Exception as e:
            print(f"❌ Error in PPT generation: {e}")
            import traceback
            traceback.print_exc()
            return None


# Example usage
if __name__ == "__main__":
    generator = StockPPTGenerator()
    
    # Test with a stock symbol
    test_symbol = "RELIANCE.NS"
    output_file = generator.generate_ppt(test_symbol)
    
    if output_file:
        print(f"✅ Success! PPT generated: {output_file}")
    else:
        print("❌ Failed to generate PPT")
